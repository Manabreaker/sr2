#!/usr/bin/env python3
"""Build a presentation for the macro time-series project."""

from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DEPS = ROOT / ".python_deps"
if DEPS.exists():
    sys.path.insert(0, str(DEPS))

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT = ROOT / "reports"
PPTX_PATH = OUT / "okun_macro_project_presentation.pptx"

COLORS = {
    "bg": RGBColor(247, 249, 251),
    "white": RGBColor(255, 255, 255),
    "ink": RGBColor(29, 36, 44),
    "muted": RGBColor(91, 105, 120),
    "accent": RGBColor(22, 84, 126),
    "accent_2": RGBColor(204, 92, 41),
    "line": RGBColor(215, 222, 229),
    "soft": RGBColor(232, 239, 245),
}

FONT = "Arial"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def img(rel: str) -> Path:
    path = OUT / rel
    if not path.exists():
        raise FileNotFoundError(path)
    return path


def set_run(run, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def add_rect(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["line"], radius=False):
    shape = slide.shapes.add_shape(1 if not radius else 5, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.6)
    return shape


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.0), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_run(run, 25, bold=True, color=COLORS["ink"])
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.86), Inches(12.0), Inches(0.32))
        stf = sub.text_frame
        stf.clear()
        stf.margin_left = 0
        p = stf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        set_run(run, 10, color=COLORS["muted"])


def add_footer(slide, num: int) -> None:
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(7.12), Inches(12.2), Pt(0.8))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.color.rgb = COLORS["line"]
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.17), Inches(9.8), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Временные ряды и кластеризация макродинамики"
    set_run(run, 8, color=COLORS["muted"])
    nbox = slide.shapes.add_textbox(Inches(12.1), Inches(7.17), Inches(0.7), Inches(0.25))
    ntf = nbox.text_frame
    ntf.clear()
    p = ntf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num:02d}"
    set_run(run, 8, color=COLORS["muted"])


def add_text(slide, text: str, x, y, w, h, size: int = 13, color=COLORS["ink"], bold=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size, bold=bold, color=color)
    return box


def add_bullets(slide, bullets: list[str], x, y, w, h, size: int = 14, title: str | None = None) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.02)
    first = True
    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_run(run, size + 1, bold=True, color=COLORS["accent"])
        first = False
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = bullet
        p.level = 0
        p.space_after = Pt(6)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["ink"]


def add_tag(slide, text: str, x, y, w, accent=COLORS["accent"]) -> None:
    shape = add_rect(slide, x, y, w, Inches(0.34), fill=COLORS["soft"], line=COLORS["soft"], radius=True)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, 9, bold=True, color=accent)


def add_metric(slide, label: str, value: str, x, y, w, h, accent=COLORS["accent"]) -> None:
    shape = add_rect(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["line"])
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = value
    set_run(run, 19, bold=True, color=accent)
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = FONT
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = COLORS["muted"]


def add_image(slide, path: Path, x, y, w, h, line=True) -> None:
    pic = slide.shapes.add_picture(str(path), x, y)
    scale = min(w / pic.width, h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(x + (w - pic.width) / 2)
    pic.top = int(y + (h - pic.height) / 2)
    if line:
        pic.line.color.rgb = COLORS["line"]
        pic.line.width = Pt(0.5)


def add_image_panel(slide, path: Path, x, y, w, h, caption: str | None = None) -> None:
    add_rect(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["line"])
    cap_h = Inches(0.32) if caption else 0
    add_image(slide, path, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), h - Inches(0.24) - cap_h)
    if caption:
        add_text(slide, caption, x + Inches(0.18), y + h - Inches(0.36), w - Inches(0.36), Inches(0.22), size=8, color=COLORS["muted"])


def add_table(slide, rows: list[list[str]], x, y, w, h, col_widths: list[float] | None = None, font_size: int = 9) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = int(w * width)
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["soft"] if r_idx == 0 else COLORS["white"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT
                paragraph.font.size = Pt(font_size)
                paragraph.font.bold = r_idx == 0
                paragraph.font.color.rgb = COLORS["ink"]


def add_takeaway(slide, text: str, x=Inches(0.72), y=Inches(6.35), w=Inches(11.9), h=Inches(0.48)) -> None:
    shape = add_rect(slide, x, y, w, h, fill=RGBColor(255, 252, 246), line=RGBColor(238, 209, 176))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.07)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, 11, bold=True, color=COLORS["accent_2"])


def blank(prs: Presentation, slide_no: int, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    add_title(slide, title, subtitle)
    add_footer(slide, slide_no)
    return slide


def cluster_rows() -> list[list[str]]:
    path = OUT / "microqualification" / "tables" / "cluster_characteristics.csv"
    df = pd.read_csv(path)
    rows = [["Кластер", "Стран", "Рост ВВП", "Безработица", "Реальная ставка"]]
    for _, row in df.sort_values("cluster").iterrows():
        name = str(row["cluster_name"]).replace("Кластер ", "")
        rows.append(
            [
                name,
                f"{int(row['n_countries'])}",
                f"{row['gdp_growth_mean']:.2f}%",
                f"{row['unemployment_mean']:.2f}%",
                f"{row['real_interest_mean']:.2f}%",
            ]
        )
    return rows


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Временные ряды и кластеризация макроэкономической динамики"
    prs.core_properties.subject = "Презентация по проекту анализа макроэкономических временных рядов"
    prs.core_properties.author = "Автор проекта"

    # 1
    slide = blank(prs, 1, "Временные ряды и кластеризация макродинамики", "Что было сделано в проекте и какие выводы получены")
    add_image(slide, img("figures/01_levels_us_macro_series.png"), Inches(6.35), Inches(1.15), Inches(6.2), Inches(4.55))
    add_bullets(
        slide,
        [
            "Основная часть: США, связь роста ВВП, безработицы, инфляции и ставки.",
            "Методы: ADF/KPSS, ACF/PACF, ADL, VAR, Granger, IRF/FEVD, VECM/ECM.",
            "Дополнение: кластеризация 22 стран по динамике макропоказателей.",
        ],
        Inches(0.72),
        Inches(1.35),
        Inches(5.5),
        Inches(2.0),
        size=15,
    )
    add_metric(slide, "квартальных наблюдений США", "145", Inches(0.72), Inches(4.0), Inches(1.7), Inches(0.92))
    add_metric(slide, "показателя FRED", "4", Inches(2.65), Inches(4.0), Inches(1.5), Inches(0.92), COLORS["accent_2"])
    add_metric(slide, "страны World Bank", "22", Inches(4.35), Inches(4.0), Inches(1.65), Inches(0.92))
    add_takeaway(slide, "Логика защиты: от подготовки рядов к моделям, затем к международной кластеризации.")

    # 2
    slide = blank(prs, 2, "1. Данные: исходные уровни", "США, FRED, 1990Q1-2026Q1")
    add_image_panel(slide, img("figures/01_levels_us_macro_series.png"), Inches(0.65), Inches(1.15), Inches(8.15), Inches(5.65))
    add_bullets(
        slide,
        [
            "ВВП взят как квартальный ряд GDPC1.",
            "Безработица, CPI и Fed Funds усреднены из месячных данных до кварталов.",
            "В уровнях заметны тренды, кризисные разрывы и разные масштабы измерения.",
        ],
        Inches(9.05),
        Inches(1.35),
        Inches(3.55),
        Inches(2.7),
        size=13,
        title="Что показываем",
    )
    add_tag(slide, "FRED", Inches(9.12), Inches(4.5), Inches(0.95))
    add_tag(slide, "BEA / BLS / Fed", Inches(10.18), Inches(4.5), Inches(1.75), COLORS["accent_2"])
    add_takeaway(slide, "Перед моделированием уровни нужно привести к сопоставимым и стационарным рядам.")

    # 3
    slide = blank(prs, 3, "2. Подготовка рядов", "Переход от уровней к динамике")
    add_image_panel(slide, img("figures/02_stationary_transformations_growth_changes.png"), Inches(0.65), Inches(1.08), Inches(8.25), Inches(5.78))
    add_bullets(
        slide,
        [
            "ВВП: квартальный темп роста в годовом выражении.",
            "Безработица: квартальное изменение уровня.",
            "CPI: инфляция; ставка: квартальное изменение.",
            "Так сохраняется экономический смысл и снижается риск ложной регрессии.",
        ],
        Inches(9.1),
        Inches(1.35),
        Inches(3.5),
        Inches(3.45),
        size=13,
        title="Преобразования",
    )
    add_takeaway(slide, "Модели ADL и VAR строятся не на исходных уровнях, а на динамических показателях.")

    # 4
    slide = blank(prs, 4, "3. Стационарность и автокорреляция", "Диагностика перед выбором модели")
    add_image_panel(slide, img("figures/04_acf_pacf_stationary_series.png"), Inches(0.65), Inches(1.05), Inches(7.6), Inches(5.85))
    add_image_panel(slide, img("calculation_screenshots/calc_01_stationarity.png"), Inches(8.45), Inches(1.05), Inches(4.05), Inches(3.0), "ADF/KPSS: расчетная проверка")
    add_bullets(
        slide,
        [
            "ADF/KPSS подтверждают необходимость преобразований.",
            "ACF/PACF помогают выбрать лаговую структуру.",
            "Для ставки видна более сильная инерционность.",
        ],
        Inches(8.55),
        Inches(4.35),
        Inches(3.95),
        Inches(1.45),
        size=12,
    )
    add_takeaway(slide, "Диагностика задает ограничения: нельзя просто регрессировать трендовые уровни друг на друга.")

    # 5
    slide = blank(prs, 5, "4. Закон Оукена: первичная связь", "Рост ВВП против изменения безработицы")
    add_image_panel(slide, img("figures/03_okun_scatter_gdp_growth_unemployment_change.png"), Inches(0.75), Inches(1.05), Inches(7.0), Inches(5.85))
    add_bullets(
        slide,
        [
            "Связь отрицательная: при ускорении выпуска безработица обычно растет слабее или снижается.",
            "COVID-кварталы дают экстремальные наблюдения.",
            "Поэтому дальше связь проверяется в ADL и VAR, а не только по scatter plot.",
        ],
        Inches(8.1),
        Inches(1.45),
        Inches(4.1),
        Inches(2.85),
        size=14,
        title="Интерпретация",
    )
    add_metric(slide, "коэффициент ADL для текущего роста ВВП", "-0.106", Inches(8.25), Inches(4.72), Inches(2.2), Inches(0.95), COLORS["accent_2"])
    add_metric(slide, "p-value", "<0.001", Inches(10.65), Inches(4.72), Inches(1.45), Inches(0.95))
    add_takeaway(slide, "Гипотеза закона Оукена подтверждается по знаку, но сила связи зависит от кризисных наблюдений.")

    # 6
    slide = blank(prs, 6, "5. ADL-модель", "Краткосрочная зависимость изменения безработицы")
    add_image_panel(slide, img("calculation_screenshots/calc_02_adl_var_models.png"), Inches(0.65), Inches(1.04), Inches(7.25), Inches(5.8), "Расчеты ADL, VAR и причинности")
    add_bullets(
        slide,
        [
            "Выбрана ADL-3: текущий и лаговые темпы роста ВВП, лаг безработицы, dummy 2020Q2.",
            "Adj.R² = 0.880, RMSE = 0.310: модель лучше базовых спецификаций.",
            "Dummy 2020Q2 значима: пандемический шок нельзя считать обычным наблюдением.",
        ],
        Inches(8.2),
        Inches(1.34),
        Inches(4.05),
        Inches(2.8),
        size=13,
        title="Что получилось",
    )
    add_table(
        slide,
        [
            ["Модель", "Adj.R²", "RMSE", "Выбор"],
            ["ADL-1", "0.738", "0.460", ""],
            ["ADL-2", "0.764", "0.435", ""],
            ["ADL-3", "0.880", "0.310", "да"],
        ],
        Inches(8.2),
        Inches(4.45),
        Inches(4.05),
        Inches(1.25),
        col_widths=[0.33, 0.22, 0.22, 0.23],
        font_size=9,
    )
    add_takeaway(slide, "ADL дает удобную оценку краткосрочной связи, но остатки требуют осторожной интерпретации.")

    # 7
    slide = blank(prs, 7, "6. VAR и причинность по Грейнджеру", "Система взаимосвязанных рядов")
    add_image_panel(slide, img("figures/06_granger_causality_network.png"), Inches(0.75), Inches(1.05), Inches(5.1), Inches(5.75))
    add_image_panel(slide, img("figures/09_model_residual_diagnostics.png"), Inches(6.1), Inches(1.05), Inches(6.25), Inches(3.45), "Диагностика остатков")
    add_bullets(
        slide,
        [
            "VAR(1) выбрана по AIC и устойчивости.",
            "Есть двусторонняя связь Грейнджера между ростом ВВП и изменением безработицы.",
            "Инфляция и ставка играют меньшую роль в уравнении безработицы.",
        ],
        Inches(6.25),
        Inches(4.78),
        Inches(5.9),
        Inches(1.25),
        size=12,
    )
    add_takeaway(slide, "VAR показывает не только связь показателей, но и направление прогнозной информативности.")

    # 8
    slide = blank(prs, 8, "7. IRF и FEVD", "Как шоки объясняют динамику безработицы")
    add_image_panel(slide, img("figures/07_var_irf_unemployment_response.png"), Inches(0.65), Inches(1.05), Inches(6.15), Inches(4.8), "Импульсный отклик Δu")
    add_image_panel(slide, img("figures/08_var_fevd_unemployment_change.png"), Inches(6.95), Inches(1.05), Inches(5.8), Inches(4.8), "Разложение дисперсии FEVD")
    add_bullets(
        slide,
        [
            "Положительный шок роста ВВП снижает прирост безработицы.",
            "На горизонтах 4-12 кварталов около 61% ошибки прогноза Δu связано с шоками выпуска.",
            "Вклад инфляции и ставки мал в выбранной спецификации.",
        ],
        Inches(0.88),
        Inches(5.96),
        Inches(11.65),
        Inches(0.7),
        size=11,
    )
    add_takeaway(slide, "IRF/FEVD превращают коэффициенты VAR в экономическую историю о реакции на шоки.")

    # 9
    slide = blank(prs, 9, "8. VECM/ECM-проверка", "Долгосрочная связь как диагностическое дополнение")
    add_image_panel(slide, img("calculation_screenshots/calc_03_vecm_ecm.png"), Inches(0.65), Inches(1.05), Inches(7.8), Inches(5.75))
    add_bullets(
        slide,
        [
            "Уровни рядов нестационарны, поэтому VECM используется как ECM-представление системы.",
            "Проверяется, есть ли долгосрочная корректировка между ВВП, безработицей, CPI и ставкой.",
            "ECM-уравнение для Δu служит диагностикой, а не заменяет основную VAR на стационарных рядах.",
        ],
        Inches(8.75),
        Inches(1.35),
        Inches(3.7),
        Inches(3.0),
        size=13,
        title="Зачем это нужно",
    )
    add_metric(slide, "p-value ECT в ECM", "0.051", Inches(8.85), Inches(4.75), Inches(1.75), Inches(0.9), COLORS["accent_2"])
    add_metric(slide, "значимая ставка", "p=0.023", Inches(10.85), Inches(4.75), Inches(1.6), Inches(0.9))
    add_takeaway(slide, "VECM/ECM закрывает вопрос о долгосрочной структуре, но основной вывод остается краткосрочным.")

    # 10
    slide = blank(prs, 10, "9. Кластеризация стран", "Сравнение не уровней, а формы динамики")
    add_image_panel(slide, img("microqualification/figures/14_cluster_average_profiles.png"), Inches(0.65), Inches(1.05), Inches(6.75), Inches(5.75), "Средние траектории кластеров")
    add_image_panel(slide, img("microqualification/figures/13_distance_heatmap.png"), Inches(7.65), Inches(1.05), Inches(4.65), Inches(3.75), "Матрица расстояний")
    add_bullets(
        slide,
        [
            "22 страны, 2010-2023, показатели World Bank.",
            "Сравнивались стандартизованные динамики: рост ВВП, безработица, инфляция, реальная ставка.",
            "Использованы dynamic Euclidean, correlation distance и summary features.",
        ],
        Inches(7.75),
        Inches(5.05),
        Inches(4.45),
        Inches(1.05),
        size=10.5,
    )
    add_takeaway(slide, "Кластеризация отвечает на вопрос: какие страны движутся похоже, даже если уровни разные.")

    # 11
    slide = blank(prs, 11, "10. Результаты кластеризации", "k-means и иерархический подход")
    add_image_panel(slide, img("microqualification/figures/11_cluster_pca_scatter.png"), Inches(0.65), Inches(1.05), Inches(5.45), Inches(4.95), "PCA-проекция k-means")
    add_image_panel(slide, img("microqualification/figures/12_hierarchical_dendrogram.png"), Inches(6.35), Inches(1.05), Inches(5.85), Inches(4.95), "Иерархическая кластеризация")
    add_bullets(
        slide,
        [
            "k-means удобно интерпретировать через центроиды и типичного представителя.",
            "Иерархическая кластеризация показывает структуру близости стран.",
            "Оба подхода дают разделение на две содержательно разные группы.",
        ],
        Inches(0.85),
        Inches(6.05),
        Inches(11.2),
        Inches(0.6),
        size=10.5,
    )
    add_takeaway(slide, "Сравнение методов снижает риск принять артефакт одного алгоритма за экономический вывод.")

    # 12
    slide = blank(prs, 12, "11. Карта и смысл кластеров", "Группы стран с похожей макродинамикой")
    add_image_panel(slide, img("microqualification/figures/10_world_map_kmeans_clusters.png"), Inches(0.65), Inches(1.0), Inches(8.0), Inches(5.65))
    add_table(slide, cluster_rows(), Inches(8.9), Inches(1.35), Inches(3.75), Inches(1.45), col_widths=[0.42, 0.14, 0.16, 0.16, 0.18], font_size=7)
    add_bullets(
        slide,
        [
            "Кластер 1: умеренные ставки и ниже безработица.",
            "Кластер 2: высокие реальные ставки и более напряженный рынок труда.",
            "Для прогноза внешних шоков страну логичнее сравнивать с ее кластером.",
        ],
        Inches(9.05),
        Inches(3.25),
        Inches(3.45),
        Inches(1.65),
        size=11.5,
    )
    add_takeaway(slide, "Карта переводит численную кластеризацию в понятную географическую интерпретацию.")

    # 13
    slide = blank(prs, 13, "12. Типичный представитель и вывод", "Перу как представитель второго кластера")
    add_image_panel(slide, img("microqualification/figures/15_representative_time_series.png"), Inches(0.65), Inches(1.05), Inches(6.75), Inches(4.95), "Макродинамика представителя")
    add_image_panel(slide, img("microqualification/figures/16_representative_okun_scatter.png"), Inches(7.65), Inches(1.05), Inches(4.6), Inches(4.95), "Оукен для представителя")
    add_bullets(
        slide,
        [
            "Для США коэффициент Оукена в ADL около -0.106; для Перу около -0.227.",
            "Знак совпадает, но сила связи и условия денежно-кредитной политики различаются.",
            "Итог: модели внутри страны и кластеризация стран дополняют друг друга.",
        ],
        Inches(0.85),
        Inches(6.05),
        Inches(11.2),
        Inches(0.62),
        size=10.5,
    )
    add_takeaway(slide, "Главный вывод: реакцию страны на шок лучше прогнозировать через ее собственную модель и похожие страны.")

    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    print(PPTX_PATH)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
